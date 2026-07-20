import json
import base64
from pathlib import Path
from tempfile import TemporaryDirectory
from unittest.mock import MagicMock, patch

import fitz
from django.conf import settings
from django.http import QueryDict
from django.test import TestCase, override_settings
from django.utils import timezone
from openpyxl import load_workbook
from docx import Document as WordDocument
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from playwright.sync_api import Error as PlaywrightError

from . import mail_backend
from .models import Job, SignatureAuditLog
from .operations import (
    _allow_browser_request,
    compress_pdf,
    flatten_edits,
    html_to_pdf,
    ocr_pdf,
    pdf_to_excel,
    pdf_to_powerpoint,
    remove_watermark_pdf,
    watermark_pdf,
    word_to_excel,
)
from .tasks import expire_stalled_jobs, process_job
from .views import read_options, sanitize_download_name


class WatermarkTests(TestCase):
    def test_read_options_keeps_watermark_formatting(self):
        expected = {
            "watermark_font": "times",
            "watermark_size": "72",
            "watermark_bold": "true",
            "watermark_italic": "true",
            "watermark_underline": "true",
            "watermark_color": "#ff0000",
            "watermark_position": "bottom-right",
            "watermark_mosaic": "true",
            "watermark_transparency": "0.5",
            "watermark_rotation": "0",
            "watermark_from_page": "1",
            "watermark_to_page": "2",
            "watermark_layer": "under",
            "watermark_mode": "image",
            "watermark_image": "data:image/png;base64,AA==",
            "watermark_image_scale": "35",
        }
        post = QueryDict(mutable=True)
        post["options"] = json.dumps(expected)
        self.assertEqual(read_options(post), expected)

    def test_text_watermark_defaults_to_every_page(self):
        with TemporaryDirectory() as directory:
            source = Path(directory) / "source.pdf"
            output = Path(directory) / "watermarked.pdf"
            document = fitz.open()
            for page_number in range(1, 4):
                page = document.new_page(width=595, height=842)
                page.insert_text((72, 72), f"PAGE {page_number}")
            document.save(source)
            document.close()

            watermark_pdf(
                [source],
                output,
                {"text": "ALL PAGES", "watermark_rotation": "0", "watermark_transparency": "1"},
            )

            with fitz.open(output) as result:
                self.assertEqual(["ALL PAGES" in page.get_text() for page in result], [True, True, True])

    def test_logo_watermark_applies_only_to_selected_pages(self):
        with TemporaryDirectory() as directory:
            source = Path(directory) / "source.pdf"
            output = Path(directory) / "watermarked.pdf"
            document = fitz.open()
            for _ in range(3):
                document.new_page(width=595, height=842)
            document.save(source)
            document.close()
            pixmap = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 80, 40), False)
            pixmap.clear_with(0x176B47)
            image_data = "data:image/png;base64," + base64.b64encode(pixmap.tobytes("png")).decode("ascii")

            watermark_pdf(
                [source],
                output,
                {
                    "watermark_mode": "image",
                    "watermark_image": image_data,
                    "watermark_from_page": "2",
                    "watermark_to_page": "3",
                    "watermark_rotation": "0",
                    "watermark_transparency": "1",
                },
            )

            with fitz.open(output) as result:
                self.assertEqual([len(page.get_images(full=True)) for page in result], [0, 1, 1])

    def test_watermark_pdf_applies_size_color_and_font(self):
        with TemporaryDirectory() as directory:
            source = Path(directory) / "source.pdf"
            output = Path(directory) / "watermarked.pdf"
            document = fitz.open()
            document.new_page(width=595, height=842)
            document.save(source)
            document.close()

            watermark_pdf(
                [source],
                output,
                {
                    "text": "FORMATTED",
                    "watermark_font": "times",
                    "watermark_size": "72",
                    "watermark_bold": "true",
                    "watermark_italic": "true",
                    "watermark_color": "#ff0000",
                    "watermark_transparency": "1",
                    "watermark_rotation": "-45",
                },
            )

            with fitz.open(output) as result:
                lines = [
                    line
                    for block in result[0].get_text("dict")["blocks"]
                    for line in block.get("lines", [])
                    if any(span.get("text") == "FORMATTED" for span in line.get("spans", []))
                ]
                spans = [span for line in lines for span in line["spans"] if span.get("text") == "FORMATTED"]
                self.assertEqual(len(spans), 1)
                self.assertAlmostEqual(spans[0]["size"], 72, places=1)
                self.assertEqual(spans[0]["color"] & 0xFFFFFF, 0xFF0000)
                self.assertIn("Times", spans[0]["font"])
                self.assertLess(lines[0]["dir"][1], 0)

    def test_removes_diagonal_watermark_from_only_one_page(self):
        with TemporaryDirectory() as directory:
            source = Path(directory) / "source.pdf"
            watermarked = Path(directory) / "watermarked.pdf"
            output = Path(directory) / "removed.pdf"
            document = fitz.open()
            for page_number in (1, 2):
                page = document.new_page(width=595, height=842)
                page.insert_text((72, 72), f"ORIGINAL PAGE {page_number}")
            document.save(source)
            document.close()

            watermark_pdf(
                [source],
                watermarked,
                {
                    "text": "CONFIDENTIAL",
                    "watermark_size": "72",
                    "watermark_rotation": "-45",
                    "watermark_from_page": "1",
                    "watermark_to_page": "1",
                },
            )
            remove_watermark_pdf([watermarked], output, {})

            with fitz.open(output) as result:
                text = "\n".join(page.get_text() for page in result)
                self.assertNotIn("CONFIDENTIAL", text)
                self.assertIn("ORIGINAL PAGE 1", text)
                self.assertIn("ORIGINAL PAGE 2", text)

    def test_removes_legacy_standalone_stream_without_redacting_content(self):
        with TemporaryDirectory() as directory:
            source = Path(directory) / "source.pdf"
            watermarked = Path(directory) / "watermarked.pdf"
            legacy = Path(directory) / "legacy.pdf"
            output = Path(directory) / "removed.pdf"
            document = fitz.open()
            page = document.new_page(width=595, height=842)
            page.insert_text((190, 421), "PRESERVE UNDER WATERMARK", fontsize=20)
            document.save(source)
            document.close()

            watermark_pdf(
                [source],
                watermarked,
                {
                    "text": "CONFIDENTIAL",
                    "watermark_size": "72",
                    "watermark_rotation": "-45",
                },
            )
            with fitz.open(watermarked) as document:
                document.xref_set_key(document[0].xref, "NFIUWatermarkStreams", "null")
                document.save(legacy)

            remove_watermark_pdf([legacy], output, {})

            with fitz.open(output) as result:
                text = result[0].get_text()
                self.assertNotIn("CONFIDENTIAL", text)
                self.assertIn("PRESERVE UNDER WATERMARK", text)


class HtmlToPdfTests(TestCase):
    def test_inline_browser_resources_are_allowed_without_network_validation(self):
        with patch("pdf_tool.operations._public_web_url") as validate:
            _allow_browser_request("data:image/png;base64,AA==")
            _allow_browser_request("blob:https://example.com/asset-id")
            _allow_browser_request("about:blank")
        validate.assert_not_called()

    def test_busy_page_still_converts_after_bounded_settle_timeout(self):
        with TemporaryDirectory() as directory:
            output = Path(directory) / "webpage.pdf"
            manager = MagicMock()
            playwright = manager.__enter__.return_value
            browser = playwright.chromium.launch.return_value
            page = browser.new_page.return_value
            response = MagicMock(ok=True, status=200)
            page.goto.return_value = response
            page.wait_for_load_state.side_effect = PlaywrightError("page remained busy")

            with (
                patch("pdf_tool.operations.sync_playwright", return_value=manager),
                patch("pdf_tool.operations._public_web_url", return_value="https://example.com/report"),
            ):
                html_to_pdf([], output, {"url": "https://example.com/report", "print_background": "true"})

            page.goto.assert_called_once_with(
                "https://example.com/report",
                wait_until="domcontentloaded",
                timeout=45000,
            )
            page.wait_for_timeout.assert_called_once_with(1000)
            page.pdf.assert_called_once()
            browser.close.assert_called_once()

    def test_browser_is_closed_when_pdf_rendering_fails(self):
        with TemporaryDirectory() as directory:
            manager = MagicMock()
            playwright = manager.__enter__.return_value
            browser = playwright.chromium.launch.return_value
            page = browser.new_page.return_value
            page.goto.return_value = MagicMock(ok=True, status=200)
            page.pdf.side_effect = RuntimeError("render failed")

            with (
                patch("pdf_tool.operations.sync_playwright", return_value=manager),
                patch("pdf_tool.operations._public_web_url", return_value="https://example.com/report"),
                self.assertRaises(RuntimeError),
            ):
                html_to_pdf([], Path(directory) / "webpage.pdf", {"url": "https://example.com/report"})

            browser.close.assert_called_once()


class CompressionTests(TestCase):
    def test_image_compression_forces_downsampling_and_jpeg_reencoding(self):
        with TemporaryDirectory() as directory:
            source = Path(directory) / "source.pdf"
            output = Path(directory) / "compressed.pdf"
            source.write_bytes(b"source-pdf" * 100)

            def create_compressed(command, check):
                self.assertTrue(check)
                output.write_bytes(b"smaller")
                return MagicMock(returncode=0)

            with (
                patch("pdf_tool.operations.resolve_binary", return_value="gs"),
                patch("pdf_tool.operations.subprocess.run", side_effect=create_compressed) as run,
            ):
                compress_pdf([source], output, {"quality": "balanced"})

            command = run.call_args.args[0]
            self.assertIn("-dColorImageDownsampleThreshold=1.0", command)
            self.assertIn("-dGrayImageDownsampleThreshold=1.0", command)
            self.assertIn("-dMonoImageDownsampleThreshold=1.0", command)
            self.assertIn("-dColorImageFilter=/DCTEncode", command)
            self.assertIn("-dGrayImageFilter=/DCTEncode", command)
            self.assertLess(output.stat().st_size, source.stat().st_size)


class OcrTests(TestCase):
    def test_ocr_uses_resolved_command_and_dependency_environment(self):
        with TemporaryDirectory() as directory:
            source = Path(directory) / "source.pdf"
            output = Path(directory) / "searchable.pdf"
            source.write_bytes(b"pdf fixture")
            completed = MagicMock(returncode=0, stderr="", stdout="")
            with (
                patch("pdf_tool.operations._ocrmypdf_command", return_value=["python", "-m", "ocrmypdf"]),
                patch("pdf_tool.operations._ocr_environment", return_value={"PATH": "ocr-tools"}),
                patch("pdf_tool.operations.subprocess.run", return_value=completed) as run,
            ):
                ocr_pdf([source], output, {"language": "eng"})

            run.assert_called_once_with(
                ["python", "-m", "ocrmypdf", "-l", "eng", "--skip-text", str(source), str(output)],
                env={"PATH": "ocr-tools"},
                capture_output=True,
                text=True,
            )


class FlattenEditsTests(TestCase):
    def test_replaces_existing_text_using_original_text_bounds(self):
        with TemporaryDirectory() as directory:
            source = Path(directory) / "source.pdf"
            output = Path(directory) / "edited.pdf"
            document = fitz.open()
            page = document.new_page(width=300, height=400)
            page.insert_text((30, 60), "THIS ORIGINAL PHRASE MUST DISAPPEAR", fontsize=12)
            document.save(source)
            document.close()

            flatten_edits(
                [source],
                output,
                {
                    "annotations": {
                        "pages": [{
                            "page": 0,
                            "objects": [{
                                "type": "text", "x": 0.1, "y": 0.1, "width": 0.4, "height": 0.05,
                                "erase": True, "erase_x": 0.09, "erase_y": 0.09,
                                "erase_width": 0.82, "erase_height": 0.09,
                                "text": "REPLACED", "font_size": 12, "viewport_width": 300,
                            }],
                        }],
                    },
                },
            )

            with fitz.open(output) as result:
                text = result[0].get_text()
                self.assertNotIn("ORIGINAL PHRASE", text)
                self.assertIn("REPLACED", text)

    def test_flattens_text_shape_and_path(self):
        with TemporaryDirectory() as directory:
            source = Path(directory) / "source.pdf"
            output = Path(directory) / "edited.pdf"
            document = fitz.open()
            document.new_page(width=300, height=400)
            document.save(source)
            document.close()
            pixmap = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 12, 12), False)
            pixmap.clear_with(0xD92D20)
            image_data = "data:image/png;base64," + base64.b64encode(pixmap.tobytes("png")).decode("ascii")

            flatten_edits(
                [source],
                output,
                {
                    "annotations": {
                        "pages": [{
                            "page": 0,
                            "objects": [
                                {"type": "text", "x": 0.1, "y": 0.1, "width": 0.5, "height": 0.1, "text": "Approved", "font_size": 18, "viewport_width": 300},
                                {"type": "highlight", "x": 0.1, "y": 0.3, "width": 0.5, "height": 0.05, "fill": "#fde047", "opacity": 0.4},
                                {"type": "path", "points": [[0.1, 0.5], [0.5, 0.6]], "color": "#d92d20", "stroke_width": 3},
                                {"type": "line", "x": 0.1, "y": 0.65, "width": 0.3, "height": 0.02, "color": "#d92d20", "stroke_width": 3},
                                {"type": "arrow", "x": 0.5, "y": 0.65, "width": 0.3, "height": 0.08, "color": "#176b47", "stroke_width": 3},
                                {"type": "stamp", "x": 0.1, "y": 0.72, "width": 0.3, "height": 0.08, "text": "APPROVED", "color": "#176b47"},
                                {"type": "image", "x": 0.5, "y": 0.72, "width": 0.15, "height": 0.1, "image": image_data},
                                {"type": "signature_field", "x": 0.1, "y": 0.84, "width": 0.4, "height": 0.08, "field_name": "ApprovalSignature"},
                            ],
                        }],
                    },
                },
            )

            with fitz.open(output) as result:
                self.assertEqual(len(result), 1)
                self.assertIn("Approved", result[0].get_text())
                self.assertIn("APPROVED", result[0].get_text())
                self.assertEqual([widget.field_name for widget in result[0].widgets()], ["ApprovalSignature"])


class PdfToPowerPointTests(TestCase):
    def test_editable_mode_retains_text_images_and_vector_graphics(self):
        with TemporaryDirectory() as directory:
            source = Path(directory) / "source.pdf"
            output = Path(directory) / "converted.pptx"
            document = fitz.open()
            page = document.new_page(width=720, height=405)
            page.draw_rect(fitz.Rect(60, 80, 300, 250), color=(0, 0, 0), fill=(0.8, 0.9, 1), width=3)
            page.draw_line((80, 300), (640, 110), color=(1, 0, 0), width=6)
            page.draw_bezier((80, 330), (220, 150), (480, 360), (640, 180), color=(0, 0.5, 0), width=4)
            page.insert_text((70, 55), "EDITABLE CHART", fontsize=24)
            pixmap = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 40, 40), False)
            pixmap.clear_with(0x3366CC)
            page.insert_image(fitz.Rect(540, 250, 620, 330), pixmap=pixmap)
            document.save(source)
            document.close()

            pdf_to_powerpoint([source], output, {"pptx_mode": "editable"})

            deck = Presentation(output)
            self.assertEqual(len(deck.slides), 1)
            shapes = list(deck.slides[0].shapes)
            shape_types = [shape.shape_type for shape in shapes]
            self.assertIn(MSO_SHAPE_TYPE.AUTO_SHAPE, shape_types)
            self.assertIn(MSO_SHAPE_TYPE.LINE, shape_types)
            self.assertIn(MSO_SHAPE_TYPE.FREEFORM, shape_types)
            self.assertIn(MSO_SHAPE_TYPE.PICTURE, shape_types)
            self.assertIn(MSO_SHAPE_TYPE.TEXT_BOX, shape_types)
            self.assertIn("EDITABLE CHART", "\n".join(shape.text for shape in shapes if shape.has_text_frame))

    def test_exact_image_mode_remains_available(self):
        with TemporaryDirectory() as directory:
            source = Path(directory) / "source.pdf"
            output = Path(directory) / "converted.pptx"
            document = fitz.open()
            page = document.new_page(width=595, height=842)
            page.insert_text((72, 72), "EXACT PAGE", fontsize=20)
            document.save(source)
            document.close()

            pdf_to_powerpoint([source], output, {"pptx_mode": "image", "dpi": "120"})

            deck = Presentation(output)
            shapes = list(deck.slides[0].shapes)
            self.assertEqual(len(shapes), 1)
            self.assertEqual(shapes[0].shape_type, MSO_SHAPE_TYPE.PICTURE)


def draw_bordered_table(page, data, x0=72, y0=140, col_width=140, row_height=24):
    for row_number in range(len(data) + 1):
        page.draw_line((x0, y0 + row_number * row_height), (x0 + len(data[0]) * col_width, y0 + row_number * row_height))
    for col_number in range(len(data[0]) + 1):
        page.draw_line((x0 + col_number * col_width, y0), (x0 + col_number * col_width, y0 + len(data) * row_height))
    for row_number, row in enumerate(data):
        for col_number, cell in enumerate(row):
            page.insert_text((x0 + col_number * col_width + 6, y0 + row_number * row_height + 16), str(cell), fontsize=10)


class PdfToExcelTests(TestCase):
    def convert(self, build_document) -> list[tuple]:
        with TemporaryDirectory() as directory:
            source = Path(directory) / "source.pdf"
            output = Path(directory) / "converted.xlsx"
            document = fitz.open()
            build_document(document)
            document.save(source)
            document.close()
            pdf_to_excel([source], output, {})
            workbook = load_workbook(output)
            self.assertEqual(workbook.sheetnames, ["Content"], "conversion must always produce exactly one sheet")
            return [row for row in workbook.active.iter_rows(values_only=True)]

    def test_sixteen_page_pdf_produces_single_sheet_with_all_pages(self):
        def build(document):
            for page_number in range(1, 17):
                page = document.new_page()
                page.insert_text((72, 72), f"Section {page_number} heading", fontsize=14)

        rows = self.convert(build)
        texts = [row[0] for row in rows if row and row[0]]
        self.assertIn("Section 1 heading", texts)
        self.assertIn("Section 16 heading", texts)

    def test_bordered_table_lands_in_separate_columns(self):
        def build(document):
            page = document.new_page()
            draw_bordered_table(page, [["Item", "Quantity", "Amount"], ["Alpha", "10", "2,500.00"]])

        rows = self.convert(build)
        self.assertIn(("Item", "Quantity", "Amount"), rows)
        self.assertIn(("Alpha", "10", "2,500.00"), rows)

    def test_text_only_pdf_yields_clean_lines_without_coordinates(self):
        def build(document):
            page = document.new_page()
            page.insert_text((72, 72), "First narrative line", fontsize=11)
            page.insert_text((72, 100), "Second narrative line", fontsize=11)

        rows = self.convert(build)
        texts = [row[0] for row in rows if row and row[0]]
        self.assertEqual(texts, ["First narrative line", "Second narrative line"])
        for row in rows:
            for cell in row:
                if cell is not None:
                    self.assertNotIsInstance(cell, float, "no layout coordinates may leak into cells")

    def test_mixed_content_keeps_reading_order(self):
        def build(document):
            page = document.new_page()
            page.insert_text((72, 72), "Report heading", fontsize=14)
            draw_bordered_table(page, [["Name", "Total"], ["Beta", "7"]])
            page.insert_text((72, 260), "Closing remarks", fontsize=11)

        rows = self.convert(build)
        flattened = [next((cell for cell in row if cell), None) for row in rows]
        flattened = [cell for cell in flattened if cell]
        self.assertEqual(
            flattened,
            ["Report heading", "Name", "Beta", "Closing remarks"],
        )

    def test_image_only_pdf_converts_without_crashing(self):
        def build(document):
            page = document.new_page()
            pixmap = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 40, 40), False)
            pixmap.clear_with(0x3366CC)
            page.insert_image(fitz.Rect(72, 72, 300, 300), pixmap=pixmap)

        rows = self.convert(build)
        self.assertTrue(all(not any(row) for row in rows), "an image-only page has no extractable text")

    def test_empty_pdf_produces_valid_workbook(self):
        rows = self.convert(lambda document: document.new_page())
        self.assertTrue(all(not any(row) for row in rows))

    def test_unicode_text_survives(self):
        def build(document):
            page = document.new_page()
            page.insert_text((72, 72), "Montant: 2 500 EUR - Reunion pleniere", fontsize=11)

        rows = self.convert(build)
        texts = [row[0] for row in rows if row and row[0]]
        self.assertEqual(texts, ["Montant: 2 500 EUR - Reunion pleniere"])

    def test_wide_table_keeps_all_columns(self):
        headers = [f"C{i}" for i in range(1, 9)]
        values = [f"v{i}" for i in range(1, 9)]

        def build(document):
            page = document.new_page(width=1200, height=400)
            draw_bordered_table(page, [headers, values], col_width=130)

        rows = self.convert(build)
        self.assertIn(tuple(headers), rows)
        self.assertIn(tuple(values), rows)


class WordToExcelTests(TestCase):
    def test_docx_paragraphs_and_tables_share_one_worksheet(self):
        with TemporaryDirectory() as directory:
            source = Path(directory) / "report.docx"
            output = Path(directory) / "converted.xlsx"
            document = WordDocument()
            document.add_paragraph("Report heading")
            table = document.add_table(rows=2, cols=2)
            table.cell(0, 0).text = "Name"
            table.cell(0, 1).text = "Score"
            table.cell(1, 0).text = "Ada"
            table.cell(1, 1).text = "95"
            document.add_paragraph("Closing remarks")
            document.save(source)

            word_to_excel([source], output, {})
            workbook = load_workbook(output)
            rows = list(workbook.active.iter_rows(values_only=True))

        self.assertEqual(workbook.sheetnames, ["Content"])
        self.assertIn(("Report heading", None), rows)
        self.assertIn(("Name", "Score"), rows)
        self.assertIn(("Ada", "95"), rows)
        self.assertIn(("Closing remarks", None), rows)


class DownloadLifecycleTests(TestCase):
    """Every uploaded document must vanish from disk the moment it has served
    its purpose: inputs when processing ends, the output when it is downloaded."""

    def setUp(self):
        self._media = TemporaryDirectory()
        self._override = self.settings(MEDIA_ROOT=Path(self._media.name))
        self._override.enable()
        self.addCleanup(self._override.disable)
        self.addCleanup(self._media.cleanup)

    def make_job(self, operation="pdf_to_excel"):
        with patch("pdf_tool.views.process_job.delay"):
            document = fitz.open()
            page = document.new_page()
            page.insert_text((72, 72), "Lifecycle test document", fontsize=12)
            pdf_bytes = document.tobytes()
            document.close()
            from django.core.files.uploadedfile import SimpleUploadedFile
            upload = SimpleUploadedFile("Board Papers.pdf", pdf_bytes, content_type="application/pdf")
            response = self.client.post("/jobs/", {"operation": operation, "files": upload})
            self.assertEqual(response.status_code, 200)
            job_id = response.json()["id"]
        process_job(job_id)
        return job_id, Path(settings.MEDIA_ROOT) / "jobs" / job_id

    def test_input_files_deleted_the_moment_processing_finishes(self):
        job_id, job_dir = self.make_job()
        self.assertFalse((job_dir / "input").exists(), "uploaded originals must not outlive processing")
        self.assertTrue((job_dir / "output").exists())

    def test_output_deleted_immediately_after_download(self):
        job_id, job_dir = self.make_job()
        response = self.client.get(f"/jobs/{job_id}/download/")
        self.assertEqual(response.status_code, 200)
        b"".join(response.streaming_content)
        response.close()
        self.assertFalse(job_dir.exists(), "nothing may remain on disk after the download completes")
        self.assertEqual(self.client.get(f"/jobs/{job_id}/download/").status_code, 404)

    def test_user_chosen_filename_is_used_for_download(self):
        job_id, job_dir = self.make_job()
        response = self.client.get(f"/jobs/{job_id}/download/", {"filename": "My Annual Returns"})
        self.assertEqual(response.headers["Content-Disposition"], 'attachment; filename="My Annual Returns.xlsx"')
        b"".join(response.streaming_content)
        response.close()

    def test_filename_sanitization(self):
        self.assertEqual(sanitize_download_name("../../etc/passwd", "r.xlsx"), "passwd.xlsx")
        self.assertEqual(sanitize_download_name('bad<>:"|?*chars', "r.xlsx"), "badchars.xlsx")
        self.assertEqual(sanitize_download_name("   ", "r.xlsx"), "r.xlsx")
        self.assertEqual(sanitize_download_name(None, "r.xlsx"), "r.xlsx")
        self.assertEqual(sanitize_download_name("evil.exe", "r.xlsx"), "evil.xlsx")

    @override_settings(EMAIL_BACKEND="pdf_tool.mail_backend.HttpMailBackend")
    def test_failed_job_also_deletes_inputs_and_emails_support(self):
        with patch("pdf_tool.mail_backend.urlopen") as post, patch("pdf_tool.views.process_job.delay"):
            post.return_value.__enter__.return_value.status = 200
            from django.core.files.uploadedfile import SimpleUploadedFile
            upload = SimpleUploadedFile("not-a-pdf.pdf", b"this is not a pdf at all", content_type="application/pdf")
            response = self.client.post("/jobs/", {"operation": "pdf_to_excel", "files": upload})
            job_id = response.json()["id"]
            process_job(job_id)
        job = Job.objects.get(id=job_id)
        self.assertEqual(job.status, Job.Status.FAILED)
        self.assertNotIn(str(settings.MEDIA_ROOT), job.error, "server paths must never leak to users")
        self.assertNotIn("\\", job.error)
        job_dir = Path(settings.MEDIA_ROOT) / "jobs" / job_id
        self.assertFalse((job_dir / "input").exists(), "inputs must be deleted even when the job fails")
        post.assert_called_once()
        request = post.call_args.args[0]
        self.assertEqual(request.full_url, settings.MAIL_SERVER_URL)
        payload = json.loads(request.data.decode("utf-8"))
        self.assertEqual(payload["recipients"], ["add@nfiu.gov.ng"])
        self.assertIn("job failed", payload["subject"])
        self.assertIn(job_id, payload["body"])


class StalledJobTests(TestCase):
    def test_stalled_queued_job_is_failed_by_watchdog(self):
        job = Job.objects.create(operation="pdf_to_excel")
        stale_time = timezone.now() - timezone.timedelta(minutes=settings.JOB_MAX_PROCESSING_MINUTES + 1)
        Job.objects.filter(pk=job.pk).update(updated_at=stale_time)

        with patch("pdf_tool.tasks.notify_support_of_failure") as notify:
            self.assertEqual(expire_stalled_jobs(), 1)

        job.refresh_from_db()
        self.assertEqual(job.status, Job.Status.FAILED)
        self.assertIn("took too long", job.error)
        notify.assert_called_once()


class SignatureAuditLogTests(TestCase):
    def test_audit_record_is_append_only(self):
        job = Job.objects.create(operation="edit")
        audit = SignatureAuditLog.objects.create(
            job=job,
            signer_name="A. Signer",
            signer_email="signer@example.com",
            ip_address="127.0.0.1",
            user_agent="test",
            signed_at=timezone.now(),
            input_sha256="a" * 64,
            output_sha256="b" * 64,
        )
        self.assertEqual(len(audit.entry_hash), 64)
        with self.assertRaises(ValueError):
            audit.save()
        with self.assertRaises(ValueError):
            audit.delete()
        with self.assertRaises(ValueError):
            SignatureAuditLog.objects.filter(pk=audit.pk).update(signer_name="Changed")
        with self.assertRaises(ValueError):
            SignatureAuditLog.objects.filter(pk=audit.pk).delete()
