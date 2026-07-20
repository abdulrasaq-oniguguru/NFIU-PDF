import base64
import hashlib
import math
import os
import ipaddress
import re
import socket
import shutil
import subprocess
import sys
import zipfile
from collections import defaultdict
from io import BytesIO
from pathlib import Path

import fitz
import pdf2docx.text.TextSpan as _pdf2docx_textspan
from openpyxl import Workbook
from openpyxl.utils import get_column_letter
from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from pdf2docx import Converter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt
from pypdf import PdfReader, PdfWriter
from pypdf.generic import ArrayObject, DictionaryObject, FloatObject, NameObject, NumberObject, TextStringObject
from playwright.sync_api import Error as PlaywrightError, sync_playwright

# PyMuPDF's rebased backend (>=1.24) can return span colors as signed integers,
# but pdf2docx's rgb_component() assumes unsigned 0-16777215 and crashes on
# negative input (hex() of a negative int keeps the sign, breaking int(..., 16)).
# Mask to 24 bits to recover the intended unsigned color before pdf2docx sees it.
_original_rgb_component = _pdf2docx_textspan.rgb_component


def _safe_rgb_component(srgb: int):
    return _original_rgb_component(srgb & 0xFFFFFF)


_pdf2docx_textspan.rgb_component = _safe_rgb_component

PDF_EXTENSIONS = {".pdf"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp"}
OFFICE_EXTENSIONS = {".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx", ".odt", ".ods", ".odp"}


def _public_web_url(value: str) -> str:
    from urllib.parse import urlparse

    parsed = urlparse((value or "").strip())
    if parsed.scheme not in {"http", "https"} or not parsed.hostname or parsed.username or parsed.password:
        raise OperationError("Enter a valid public http or https URL")
    try:
        addresses = socket.getaddrinfo(parsed.hostname, parsed.port or (443 if parsed.scheme == "https" else 80))
    except socket.gaierror as exc:
        raise OperationError("The website address could not be resolved") from exc
    for address in addresses:
        ip = ipaddress.ip_address(address[4][0])
        if not ip.is_global:
            raise OperationError("Private, local, and reserved network addresses are not allowed")
    return parsed.geturl()


def _allow_browser_request(value: str) -> None:
    from urllib.parse import urlparse

    # These schemes are generated within the already-approved page and cannot
    # initiate a new request to an internal network address.
    if urlparse(value).scheme in {"data", "blob", "about"}:
        return
    _public_web_url(value)


def html_to_pdf(files: list[Path], output: Path, options: dict) -> Path:
    url = _public_web_url(str(options.get("url") or ""))
    width = min(max(int(options.get("screen_width") or 1440), 320), 2560)
    page_size = str(options.get("page_size") or "A4")
    if page_size not in {"A4", "Letter", "Legal"}:
        page_size = "A4"
    landscape = str(options.get("orientation") or "portrait") == "landscape"
    long_page = _truthy(options.get("one_long_page"))
    with sync_playwright() as playwright:
        browser = None
        launch_errors = []
        for launch_options in ({"headless": True}, {"headless": True, "channel": "chrome"}, {"headless": True, "channel": "msedge"}):
            try:
                browser = playwright.chromium.launch(**launch_options)
                break
            except PlaywrightError as exc:
                launch_errors.append(str(exc))
        if browser is None:
            raise OperationError("No Chromium browser is available. Run: playwright install chromium")
        try:
            page = browser.new_page(viewport={"width": width, "height": 900}, device_scale_factor=1)

            def guard(route):
                try:
                    _allow_browser_request(route.request.url)
                    route.continue_()
                except OperationError:
                    route.abort()

            page.route("**/*", guard)
            # The load event can be held open indefinitely by advertisements or
            # broken third-party resources. DOMContentLoaded gives a stable document,
            # after which we allow a bounded network-settle period for images/styles.
            try:
                response = page.goto(url, wait_until="domcontentloaded", timeout=45000)
            except PlaywrightError as exc:
                raise OperationError("The webpage could not be loaded") from exc
            if not response or not response.ok:
                status = response.status if response else "no response"
                raise OperationError(f"The webpage could not be loaded (server returned {status})")
            try:
                page.wait_for_load_state("networkidle", timeout=8000)
            except PlaywrightError:
                # Busy pages may never become idle; give already-started rendering a
                # final brief window without turning that into a conversion failure.
                page.wait_for_timeout(1000)
            page.emulate_media(media="screen")
            pdf_options = {
                "path": str(output),
                "print_background": _truthy(options.get("print_background", True)),
                "landscape": landscape,
                "margin": {"top": "10mm", "right": "10mm", "bottom": "10mm", "left": "10mm"},
            }
            if long_page:
                height = min(max(page.evaluate("Math.max(document.body.scrollHeight, document.documentElement.scrollHeight)"), 1), 20000)
                pdf_options.update({"width": f"{width}px", "height": f"{height}px", "margin": {"top": "0", "right": "0", "bottom": "0", "left": "0"}})
            else:
                pdf_options["format"] = page_size
            page.pdf(**pdf_options)
        finally:
            browser.close()
    return output


class OperationError(RuntimeError):
    pass


def run_operation(operation: str, files: list[Path], output_dir: Path, options: dict) -> Path:
    if operation not in OPERATIONS:
        raise OperationError("Unknown operation")
    output_name, handler = OPERATIONS[operation]
    output = output_dir / output_name
    return handler(files, output, options)


def merge_pdfs(files: list[Path], output: Path, options: dict) -> Path:
    writer = PdfWriter()
    for file_path in files:
        ensure_pdf(file_path)
        reader = PdfReader(str(file_path))
        for page in reader.pages:
            writer.add_page(page)
    with output.open("wb") as handle:
        writer.write(handle)
    return output


def split_pdf(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    reader = PdfReader(str(source))
    every = max(int(options.get("every") or 1), 1)
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for index in range(0, len(reader.pages), every):
            writer = PdfWriter()
            for page in reader.pages[index : index + every]:
                writer.add_page(page)
            part_name = f"{source.stem}_pages_{index + 1}-{min(index + every, len(reader.pages))}.pdf"
            part_path = output.parent / part_name
            with part_path.open("wb") as handle:
                writer.write(handle)
            archive.write(part_path, part_name)
            part_path.unlink(missing_ok=True)
    return output


def delete_pages(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    reader = PdfReader(str(source))
    delete_set = set(parse_page_spec(options.get("pages") or "", len(reader.pages)))
    if not delete_set:
        raise OperationError("Enter pages to delete")
    writer = PdfWriter()
    for index, page in enumerate(reader.pages, start=1):
        if index not in delete_set:
            writer.add_page(page)
    if not writer.pages:
        raise OperationError("Cannot delete every page")
    with output.open("wb") as handle:
        writer.write(handle)
    return output


def reorder_pages(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    reader = PdfReader(str(source))
    page_count = len(reader.pages)
    order = parse_page_spec(options.get("pages") or "", page_count)
    if not order:
        raise OperationError("Enter the new page order")
    if str(options.get("keep_remaining") or "").lower() in ("1", "true", "yes", "on"):
        moved = set(order)
        order = order + [page for page in range(1, page_count + 1) if page not in moved]
    writer = PdfWriter()
    for page_number in order:
        writer.add_page(reader.pages[page_number - 1])
    with output.open("wb") as handle:
        writer.write(handle)
    return output


def rotate_pdf(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    degrees = int(options.get("degrees") or 90)
    reader = PdfReader(str(source))
    writer = PdfWriter()
    for page in reader.pages:
        page.rotate(degrees)
        writer.add_page(page)
    with output.open("wb") as handle:
        writer.write(handle)
    return output


def crop_pdf(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    reader = PdfReader(str(source))
    writer = PdfWriter()
    crop = options.get("crop") or {}
    use_selection = all(key in crop for key in ("x", "y", "width", "height"))
    margin = max(float(options.get("margin") or 0), 0)
    for index, page in enumerate(reader.pages, start=1):
        applies = crop.get("scope", "all") == "all" or index == int(crop.get("page") or 1)
        if use_selection and applies:
            left, bottom, right, top = map(float, (page.cropbox.left, page.cropbox.bottom, page.cropbox.right, page.cropbox.top))
            page_width, page_height = right - left, top - bottom
            x = min(max(float(crop["x"]), 0), 1)
            y = min(max(float(crop["y"]), 0), 1)
            width = min(max(float(crop["width"]), 0.001), 1 - x)
            height = min(max(float(crop["height"]), 0.001), 1 - y)
            page.cropbox.lower_left = (left + x * page_width, top - (y + height) * page_height)
            page.cropbox.upper_right = (left + (x + width) * page_width, top - y * page_height)
        elif not use_selection and margin:
            page.cropbox.lower_left = (float(page.cropbox.left) + margin, float(page.cropbox.bottom) + margin)
            page.cropbox.upper_right = (float(page.cropbox.right) - margin, float(page.cropbox.top) - margin)
        writer.add_page(page)
    with output.open("wb") as handle:
        writer.write(handle)
    return output


# Named Ghostscript presets (/ebook, /printer, ...) only downsample images that
# already exceed their built-in default resolution, so a document with images
# at or below that resolution just gets re-encoded with no size benefit (and
# sometimes a net increase from added font/structure overhead). Forcing an
# explicit resolution per tier makes each tier actually downsample.
COMPRESSION_PROFILES = {
    "high": {"preset": "screen", "dpi": 96},
    "balanced": {"preset": "ebook", "dpi": 150},
    "light": {"preset": "printer", "dpi": 300},
}


def compress_pdf(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    ghostscript = resolve_binary("gs", "gswin64c", "gswin32c")
    quality = options.get("quality") or "balanced"
    profile = COMPRESSION_PROFILES.get(quality, COMPRESSION_PROFILES["balanced"])
    subprocess.run(
        [
            ghostscript,
            "-sDEVICE=pdfwrite",
            "-dCompatibilityLevel=1.4",
            f"-dPDFSETTINGS=/{profile['preset']}",
            "-dNOPAUSE",
            "-dQUIET",
            "-dBATCH",
            "-dDownsampleColorImages=true",
            f"-dColorImageResolution={profile['dpi']}",
            # Ghostscript only downsamples when an image's resolution exceeds
            # (target * threshold); the default threshold of 1.5 leaves moderately
            # over-resolution scans untouched. A threshold of 1.0 downsamples any
            # image above the target, which is what actually shrinks image PDFs.
            "-dColorImageDownsampleThreshold=1.0",
            "-dDownsampleGrayImages=true",
            f"-dGrayImageResolution={profile['dpi']}",
            "-dGrayImageDownsampleThreshold=1.0",
            "-dDownsampleMonoImages=true",
            f"-dMonoImageResolution={profile['dpi']}",
            "-dMonoImageDownsampleThreshold=1.0",
            # Re-encode color/gray images as JPEG so already-embedded images are
            # recompressed (not just passed through), improving the ratio.
            "-dAutoFilterColorImages=false",
            "-dColorImageFilter=/DCTEncode",
            "-dAutoFilterGrayImages=false",
            "-dGrayImageFilter=/DCTEncode",
            f"-sOutputFile={output}",
            str(source),
        ],
        check=True,
    )
    if output.stat().st_size >= source.stat().st_size:
        shutil.copyfile(source, output)
    return output


def protect_pdf(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    password = options.get("password")
    if not password:
        raise OperationError("Password is required")
    reader = PdfReader(str(source))
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(user_password=password, owner_password=password, algorithm="AES-256")
    with output.open("wb") as handle:
        writer.write(handle)
    return output


def unlock_pdf(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    password = options.get("password") or ""
    reader = PdfReader(str(source))
    if reader.is_encrypted and not reader.decrypt(password):
        raise OperationError("Invalid or missing PDF password")
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    with output.open("wb") as handle:
        writer.write(handle)
    return output


WATERMARK_POSITIONS = {
    "top-left": (0.2, 0.18),
    "top-center": (0.5, 0.18),
    "top-right": (0.8, 0.18),
    "middle-left": (0.2, 0.5),
    "center": (0.5, 0.5),
    "middle-right": (0.8, 0.5),
    "bottom-left": (0.2, 0.82),
    "bottom-center": (0.5, 0.82),
    "bottom-right": (0.8, 0.82),
}


def _clamp_float(value, default: float, minimum: float, maximum: float) -> float:
    try:
        number = float(value)
    except (TypeError, ValueError):
        number = default
    return min(max(number, minimum), maximum)


def _clamp_int(value, default: int, minimum: int, maximum: int) -> int:
    try:
        number = int(value)
    except (TypeError, ValueError):
        number = default
    return min(max(number, minimum), maximum)


def _watermark_font(options: dict) -> str:
    bold = _truthy(options.get("watermark_bold"))
    italic = _truthy(options.get("watermark_italic"))
    return PAGE_NUMBER_FONTS.get((options.get("watermark_font") or "helvetica", bold, italic), "helv")


def _watermark_points(page, options: dict) -> list[fitz.Point]:
    rect = page.rect
    if _truthy(options.get("watermark_mosaic")):
        return [fitz.Point(rect.width * x, rect.height * y) for y in (0.22, 0.5, 0.78) for x in (0.22, 0.5, 0.78)]
    x_ratio, y_ratio = WATERMARK_POSITIONS.get(options.get("watermark_position") or "center", WATERMARK_POSITIONS["center"])
    return [fitz.Point(rect.width * x_ratio, rect.height * y_ratio)]


def _draw_watermark_text(page, text: str, point: fitz.Point, options: dict) -> None:
    fontsize = _clamp_float(options.get("watermark_size"), 48, 8, 180)
    fontname = _watermark_font(options)
    color = _hex_to_rgb(options.get("watermark_color") or "#727272")
    opacity = _clamp_float(options.get("watermark_transparency"), 0.18, 0.05, 1)
    rotation = _clamp_float(options.get("watermark_rotation"), 45, -180, 180)
    overlay = (options.get("watermark_layer") or "over") != "under"

    text_width = max(fitz.get_text_length(text, fontname=fontname, fontsize=fontsize), fontsize * 2)
    box_width = min(max(text_width + fontsize * 2, fontsize * 4), page.rect.width * 1.8)
    box_height = fontsize * 2.2
    box = fitz.Rect(point.x - box_width / 2, point.y - box_height / 2, point.x + box_width / 2, point.y + box_height / 2)
    # PyMuPDF's page-coordinate rotation is visually opposite to CSS rotation.
    morph = (point, fitz.Matrix(1, 1).prerotate(-rotation)) if rotation else None
    page.insert_textbox(
        box,
        text,
        fontsize=fontsize,
        fontname=fontname,
        color=color,
        align=fitz.TEXT_ALIGN_CENTER,
        morph=morph,
        fill_opacity=opacity,
        overlay=overlay,
    )
    if _truthy(options.get("watermark_underline")):
        angle = math.radians(rotation)
        direction = fitz.Point(math.cos(angle), math.sin(angle))
        normal = fitz.Point(-math.sin(angle), math.cos(angle))
        underline_center = point + normal * (fontsize * 0.58)
        half_width = text_width / 2
        page.draw_line(
            underline_center - direction * half_width,
            underline_center + direction * half_width,
            color=color,
            width=max(fontsize * 0.05, 0.6),
            stroke_opacity=opacity,
            overlay=overlay,
        )


def _load_watermark_image(options: dict) -> bytes:
    raw = str(options.get("watermark_image") or "")
    data = decode_data_url(raw)
    if not data:
        raise OperationError("Upload a logo or image to use as the watermark")
    return data


def _draw_watermark_image(page, image: bytes, point: fitz.Point, options: dict) -> None:
    from io import BytesIO

    from PIL import Image

    opacity = _clamp_float(options.get("watermark_transparency"), 0.18, 0.05, 1)
    rotation = _clamp_float(options.get("watermark_rotation"), 45, -180, 180)
    scale = _clamp_float(options.get("watermark_image_scale"), 30, 5, 100) / 100
    overlay = (options.get("watermark_layer") or "over") != "under"

    picture = Image.open(BytesIO(image)).convert("RGBA")
    if opacity < 1:
        alpha = picture.getchannel("A").point(lambda value: int(value * opacity))
        picture.putalpha(alpha)
    if rotation:
        # PIL rotates counter-clockwise; negate so positive angles match the text path.
        picture = picture.rotate(-rotation, expand=True, resample=Image.BICUBIC)

    target_width = page.rect.width * scale
    target_height = target_width * (picture.height / picture.width)
    box = fitz.Rect(
        point.x - target_width / 2,
        point.y - target_height / 2,
        point.x + target_width / 2,
        point.y + target_height / 2,
    )
    buffer = BytesIO()
    picture.save(buffer, format="PNG")
    page.insert_image(box, stream=buffer.getvalue(), keep_proportion=True, overlay=overlay)


def watermark_pdf(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    image_mode = (options.get("watermark_mode") or "text") == "image"
    if image_mode:
        image_bytes = _load_watermark_image(options)
    else:
        text = str(options.get("text") or "CONFIDENTIAL").strip()
        if not text:
            raise OperationError("Enter watermark text")
    document = fitz.open(source)
    try:
        first_page = _clamp_int(options.get("watermark_from_page"), 1, 1, len(document))
        last_page = _clamp_int(options.get("watermark_to_page"), len(document), first_page, len(document))
        for page_index in range(first_page - 1, last_page):
            page = document[page_index]
            original_streams = set(page.get_contents())
            for point in _watermark_points(page, options):
                if image_mode:
                    _draw_watermark_image(page, image_bytes, point, options)
                else:
                    _draw_watermark_text(page, text, point, options)
            watermark_streams = [xref for xref in page.get_contents() if xref not in original_streams]
            if watermark_streams:
                references = " ".join(f"{xref} 0 R" for xref in watermark_streams)
                document.xref_set_key(page.xref, "NFIUWatermarkStreams", f"[{references}]")
        document.save(output)
    finally:
        document.close()
    return output


def remove_watermark_pdf(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    document = fitz.open(source)
    page_count = len(document)
    threshold = max(1, math.ceil(page_count * 0.6))

    text_groups: dict[tuple[str, int], list[tuple[int, fitz.Quad]]] = defaultdict(list)
    text_group_values: dict[tuple[str, int], str] = {}
    diagonal_text_groups: set[tuple[str, int]] = set()
    for page_index in range(page_count):
        page = document[page_index]
        min_dim = min(page.rect.width, page.rect.height)
        raw = page.get_text("rawdict")
        for block in raw.get("blocks", []):
            for line in block.get("lines", []):
                direction = line.get("dir", (1, 0))
                angle = math.degrees(math.atan2(direction[1], direction[0]))
                angle_mod = abs(angle) % 90
                is_diagonal = min(angle_mod, 90 - angle_mod) > 5
                for span in line.get("spans", []):
                    chars = span.get("chars") or []
                    text = "".join(c.get("c", "") for c in chars).strip()
                    if not text:
                        continue
                    is_large = span.get("size", 0) > 0.05 * min_dim
                    if is_diagonal or is_large:
                        key = (text.lower(), round(angle))
                        text_group_values.setdefault(key, text)
                        if is_diagonal:
                            diagonal_text_groups.add(key)
                        for quad in _char_quads(direction, span, chars):
                            text_groups[key].append((page_index, quad))

    image_groups: dict[str, list[tuple[int, fitz.Rect]]] = defaultdict(list)
    for page_index in range(page_count):
        page = document[page_index]
        for image_info in page.get_images(full=True):
            xref = image_info[0]
            try:
                image_bytes = document.extract_image(xref)["image"]
            except Exception:
                continue
            digest = hashlib.md5(image_bytes).hexdigest()
            for rect in page.get_image_rects(xref):
                image_groups[digest].append((page_index, rect))

    stream_removals: dict[int, set[int]] = defaultdict(set)
    for page_index, page in enumerate(document):
        key_type, key_value = document.xref_get_key(page.xref, "NFIUWatermarkStreams")
        if key_type == "array":
            stream_removals[page_index].update(int(value) for value in re.findall(r"(\d+)\s+0\s+R", key_value))

    removed = sum(len(xrefs) for xrefs in stream_removals.values())
    for key, occurrences in text_groups.items():
        occurrence_pages = {page_index for page_index, _ in occurrences}
        if key not in diagonal_text_groups and len(occurrence_pages) < threshold:
            continue

        encoded_text = text_group_values[key].encode("latin-1", errors="ignore").hex().encode("ascii")
        stream_matches: dict[int, set[int]] = defaultdict(set)
        if encoded_text:
            for page_index in occurrence_pages:
                for xref in document[page_index].get_contents():
                    stream = document.xref_stream(xref)
                    if len(stream) <= 16384 and encoded_text in stream.lower():
                        stream_matches[page_index].add(xref)

        for page_index, quad in occurrences:
            matches = stream_matches.get(page_index)
            if matches:
                before = len(stream_removals[page_index])
                stream_removals[page_index].update(matches)
                removed += len(stream_removals[page_index]) - before
            else:
                document[page_index].add_redact_annot(quad)
                removed += 1
    for occurrences in image_groups.values():
        if len({page_index for page_index, _ in occurrences}) >= threshold:
            for page_index, rect in occurrences:
                document[page_index].add_redact_annot(rect)
                removed += 1

    if removed == 0:
        document.close()
        raise OperationError("No watermark-like content was detected")

    for page_index, xrefs in stream_removals.items():
        page = document[page_index]
        remaining = [xref for xref in page.get_contents() if xref not in xrefs]
        references = " ".join(f"{xref} 0 R" for xref in remaining)
        document.xref_set_key(page.xref, "Contents", f"[{references}]")
        document.xref_set_key(page.xref, "NFIUWatermarkStreams", "null")

    for page in document:
        page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_REMOVE)

    document.save(output, garbage=4, deflate=True)
    document.close()
    return output


def _char_quads(direction, span: dict, chars: list[dict]) -> list[fitz.Quad]:
    """Build tight per-character rotated quads from character origins.

    MuPDF reports span['bbox'] as an axis-aligned box in device space, which
    for rotated (e.g. diagonal watermark) text becomes far larger than the
    visible glyphs and can bleed into unrelated content. Redacting one quad
    per character (rather than one big quad for the whole run) keeps the
    erased area as close as possible to the actual glyph strokes, minimizing
    collateral damage where the watermark visually crosses other text.
    """
    dx, dy = direction
    norm = math.hypot(dx, dy) or 1.0
    dx, dy = dx / norm, dy / norm
    px, py = -dy, dx
    size = span.get("size", 10)
    ascender = span.get("ascender", 0.8) * size
    descender = span.get("descender", -0.2) * size
    up = fitz.Point(px, py)
    step = fitz.Point(dx, dy) * (size * 0.6)

    quads = []
    for index, char in enumerate(chars):
        start = fitz.Point(char["origin"])
        end = fitz.Point(chars[index + 1]["origin"]) if index + 1 < len(chars) else start + step
        top_left = start + up * ascender
        bottom_left = start + up * descender
        top_right = end + up * ascender
        bottom_right = end + up * descender
        quads.append(fitz.Quad(top_left, top_right, bottom_left, bottom_right))
    return quads


PAGE_NUMBER_POSITIONS = {
    "top-left": ("top", fitz.TEXT_ALIGN_LEFT, "left"),
    "top-center": ("top", fitz.TEXT_ALIGN_CENTER, "center"),
    "top-right": ("top", fitz.TEXT_ALIGN_RIGHT, "right"),
    "bottom-left": ("bottom", fitz.TEXT_ALIGN_LEFT, "left"),
    "bottom-center": ("bottom", fitz.TEXT_ALIGN_CENTER, "center"),
    "bottom-right": ("bottom", fitz.TEXT_ALIGN_RIGHT, "right"),
}

PAGE_NUMBER_MARGINS = {"small": 10, "recommended": 20, "big": 36}

PAGE_NUMBER_FONTS = {
    ("helvetica", False, False): "helv",
    ("helvetica", True, False): "hebo",
    ("helvetica", False, True): "heit",
    ("helvetica", True, True): "hebi",
    ("times", False, False): "tiro",
    ("times", True, False): "tibo",
    ("times", False, True): "tiit",
    ("times", True, True): "tibi",
    ("courier", False, False): "cour",
    ("courier", True, False): "cobo",
    ("courier", False, True): "coit",
    ("courier", True, True): "cobi",
}


def _truthy(value) -> bool:
    return str(value or "false").lower() in ("1", "true", "yes", "on")


def _hex_to_rgb(value: str) -> tuple:
    value = (value or "#1a1f2b").lstrip("#")
    if len(value) != 6:
        value = "1a1f2b"
    return tuple(int(value[i : i + 2], 16) / 255 for i in (0, 2, 4))


def page_numbers_pdf(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    document = fitz.open(source)
    start = max(int(options.get("page_number_start") or 1), 1)
    last_number = start + len(document) - 1
    template = options.get("page_number_format") or "number_total"
    custom_text = options.get("page_number_custom") or "Page {n} of {p}"
    vertical, align, horizontal = PAGE_NUMBER_POSITIONS.get(
        options.get("page_number_position") or "bottom-center", PAGE_NUMBER_POSITIONS["bottom-center"]
    )
    facing = (options.get("page_number_mode") or "single") == "facing"
    margin = PAGE_NUMBER_MARGINS.get(options.get("page_number_margin") or "recommended", 20)
    box_height = 24

    bold = _truthy(options.get("page_number_bold"))
    italic = _truthy(options.get("page_number_italic"))
    underline = _truthy(options.get("page_number_underline"))
    fontname = PAGE_NUMBER_FONTS.get((options.get("page_number_font") or "helvetica", bold, italic), "helv")
    fontsize = max(float(options.get("page_number_size") or 10), 4)
    color = _hex_to_rgb(options.get("page_number_color"))

    for index, page in enumerate(document):
        number = start + index
        if template == "page_n":
            text = f"Page {number}"
        elif template == "page_of":
            text = f"Page {number} of {last_number}"
        elif template == "custom":
            text = custom_text.replace("{n}", str(number)).replace("{p}", str(last_number))
        elif template == "number":
            text = str(number)
        else:
            text = f"{number} / {last_number}"

        page_align = align
        # On facing pages, mirror left/right placement on even pages so the
        # number always sits on the outer edge, like a printed book spread.
        if facing and index % 2 == 1 and horizontal in ("left", "right"):
            page_align = fitz.TEXT_ALIGN_RIGHT if horizontal == "left" else fitz.TEXT_ALIGN_LEFT

        rect = page.rect
        y0 = margin if vertical == "top" else rect.height - margin - box_height
        box = fitz.Rect(12, y0, rect.width - 12, y0 + box_height)
        page.insert_textbox(box, text, fontsize=fontsize, fontname=fontname, color=color, align=page_align)

        if underline:
            text_width = fitz.get_text_length(text, fontname=fontname, fontsize=fontsize)
            if page_align == fitz.TEXT_ALIGN_LEFT:
                x0 = box.x0 + 2
            elif page_align == fitz.TEXT_ALIGN_RIGHT:
                x0 = box.x1 - 2 - text_width
            else:
                x0 = box.x0 + (box.width - text_width) / 2
            underline_y = y0 + box_height - 6
            page.draw_line(
                fitz.Point(x0, underline_y),
                fitz.Point(x0 + text_width, underline_y),
                color=color,
                width=max(fontsize * 0.06, 0.6),
            )
    document.save(output)
    document.close()
    return output


def images_to_pdf(files: list[Path], output: Path, options: dict) -> Path:
    image_files = [path for path in files if path.suffix.lower() in IMAGE_EXTENSIONS]
    if not image_files:
        raise OperationError("Upload at least one image")
    document = fitz.open()
    for image_path in image_files:
        image_document = fitz.open(image_path)
        image_pdf = fitz.open("pdf", image_document.convert_to_pdf())
        document.insert_pdf(image_pdf)
        image_pdf.close()
        image_document.close()
    document.save(output)
    document.close()
    return output


def iter_pdf_images(document: "fitz.Document"):
    """Yield every embedded image in a PDF as (id, page, index, ext, width, height, data).

    ``id`` (``"<page>-<index>"``, 1-based) is stable across calls for the same
    document so the extract-images preview and the actual extraction job agree
    on which image is which.
    """
    for page_index in range(len(document)):
        page = document[page_index]
        for image_index, image_info in enumerate(page.get_images(full=True), start=1):
            xref = image_info[0]
            image = document.extract_image(xref)
            yield {
                "id": f"{page_index + 1}-{image_index}",
                "page": page_index + 1,
                "index": image_index,
                "ext": image.get("ext", "png"),
                "width": image.get("width", 0),
                "height": image.get("height", 0),
                "data": image["image"],
            }


def build_image_thumbnail(data: bytes, max_size: int = 240) -> str | None:
    """Render an embedded image's bytes as a small base64 PNG data URL for previews."""
    try:
        pixmap = fitz.Pixmap(data)
        if pixmap.colorspace is None or pixmap.colorspace.n > 3:
            pixmap = fitz.Pixmap(fitz.csRGB, pixmap)
        scale = min(1.0, max_size / max(pixmap.width, pixmap.height, 1))
        if scale < 1.0:
            new_width = max(1, round(pixmap.width * scale))
            new_height = max(1, round(pixmap.height * scale))
            # PyMuPDF's 3-arg Pixmap(src, w, h) scaled-copy path is broken in some
            # installed versions (unconditionally unpacks args as if 4 were given),
            # so always pass the source's own irect as an explicit clip.
            pixmap = fitz.Pixmap(pixmap, float(new_width), float(new_height), pixmap.irect)
        thumbnail_bytes = pixmap.tobytes("png")
    except Exception:
        return None
    return "data:image/png;base64," + base64.b64encode(thumbnail_bytes).decode("ascii")


def extract_images(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    document = fitz.open(source)
    selected = options.get("selected_images")
    selected_ids = set(selected) if isinstance(selected, list) else None
    if selected_ids is not None and not selected_ids:
        document.close()
        raise OperationError("Select at least one image to extract")
    extracted = 0
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for entry in iter_pdf_images(document):
            if selected_ids is not None and entry["id"] not in selected_ids:
                continue
            image_name = f"{source.stem}_page_{entry['page']}_image_{entry['index']}.{entry['ext']}"
            image_path = output.parent / image_name
            image_path.write_bytes(entry["data"])
            archive.write(image_path, image_name)
            image_path.unlink(missing_ok=True)
            extracted += 1
    document.close()
    if extracted == 0:
        raise OperationError("No embedded images found")
    return output


def pdf_to_images(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    dpi = int(options.get("dpi") or 200)
    document = fitz.open(source)
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for index, page in enumerate(document, start=1):
            image_name = f"{source.stem}_page_{index}.png"
            image_path = output.parent / image_name
            page.get_pixmap(dpi=dpi).save(image_path)
            archive.write(image_path, image_name)
            image_path.unlink(missing_ok=True)
    document.close()
    return output


def pdf_to_word(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    converter = Converter(str(source))
    try:
        converter.convert(str(output))
    finally:
        converter.close()
    return output


def pdf_to_excel(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    document = fitz.open(source)
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Content"
    max_columns = 1
    for page in document:
        # Real tables become real rows/columns; remaining text becomes one
        # cell per line, everything appended to the same single sheet in
        # reading order.
        tables = page.find_tables().tables
        table_areas = [fitz.Rect(table.bbox) for table in tables]
        items = [(table.bbox[1], "table", table) for table in tables]
        for block in page.get_text("blocks"):
            text = str(block[4]).strip()
            if not text:
                continue
            rect = fitz.Rect(block[:4])
            center = fitz.Point((rect.x0 + rect.x1) / 2, (rect.y0 + rect.y1) / 2)
            if any(area.contains(center) for area in table_areas):
                continue
            items.append((block[1], "text", text))
        items.sort(key=lambda item: item[0])
        for _, kind, payload in items:
            if kind == "table":
                for table_row in payload.extract():
                    cells = [str(cell).strip() if cell is not None else "" for cell in table_row]
                    if any(cells):
                        sheet.append(cells)
                        max_columns = max(max_columns, len(cells))
            else:
                for line in payload.splitlines():
                    line = line.strip()
                    if line:
                        sheet.append([line])
        sheet.append([])
    document.close()
    for column_index in range(1, max_columns + 1):
        sheet.column_dimensions[get_column_letter(column_index)].width = 40
    workbook.save(output)
    return output


def word_to_excel(files: list[Path], output: Path, options: dict) -> Path:
    """Extract Word paragraphs and tables into a single Excel worksheet."""
    source = single_file(files)
    if source.suffix.lower() != ".docx":
        raise OperationError("Word to Excel supports .docx files. Save legacy .doc files as .docx first.")

    document = Document(source)
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Content"
    max_columns = 1
    for element in document.element.body.iterchildren():
        if element.tag.endswith("}p"):
            text = Paragraph(element, document).text.strip()
            if text:
                sheet.append([text])
        elif element.tag.endswith("}tbl"):
            table = Table(element, document)
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    sheet.append(cells)
                    max_columns = max(max_columns, len(cells))
            sheet.append([])
    for column_index in range(1, max_columns + 1):
        sheet.column_dimensions[get_column_letter(column_index)].width = 40
    workbook.save(output)
    return output


EMU_PER_POINT = 12700  # 914400 EMU per inch / 72 points per inch


def _pdf_int_color_to_rgb(color: int) -> "RGBColor":
    color &= 0xFFFFFF
    return RGBColor((color >> 16) & 0xFF, (color >> 8) & 0xFF, color & 0xFF)


def _pdf_color_to_rgb(color) -> "RGBColor":
    values = tuple(color or (0, 0, 0))[:3]
    return RGBColor(*(round(min(max(float(value), 0), 1) * 255) for value in values))


def _set_powerpoint_opacity(color_elements, opacity) -> None:
    if opacity is None or float(opacity) >= 1:
        return
    alpha_value = str(round(min(max(float(opacity), 0), 1) * 100000))
    for color_element in color_elements:
        for child in list(color_element):
            if child.tag.endswith("}alpha"):
                color_element.remove(child)
        alpha = OxmlElement("a:alpha")
        alpha.set("val", alpha_value)
        color_element.append(alpha)


def _style_powerpoint_shape(shape, drawing: dict, scale: float, include_fill: bool = True) -> None:
    stroke = drawing.get("color")
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = _pdf_color_to_rgb(stroke)
        shape.line.width = Pt(max(float(drawing.get("width") or 1) * scale, 0.25))
        _set_powerpoint_opacity(
            shape._element.spPr.xpath("./a:ln/a:solidFill/a:srgbClr"),
            drawing.get("stroke_opacity"),
        )
        dash_match = re.search(r"\[([^]]*)\]", str(drawing.get("dashes") or ""))
        dash_numbers = [float(value) for value in re.findall(r"[0-9.]+", dash_match.group(1))] if dash_match else []
        if dash_numbers:
            shape.line.dash_style = (
                MSO_LINE_DASH_STYLE.ROUND_DOT
                if max(dash_numbers) <= 2
                else MSO_LINE_DASH_STYLE.DASH_DOT
                if len(dash_numbers) >= 4
                else MSO_LINE_DASH_STYLE.DASH
            )

    if not include_fill:
        return
    fill = drawing.get("fill")
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _pdf_color_to_rgb(fill)
        _set_powerpoint_opacity(
            shape._element.spPr.xpath("./a:solidFill/a:srgbClr"),
            drawing.get("fill_opacity"),
        )


def _points_match(first, second, tolerance: float = 0.1) -> bool:
    return abs(first.x - second.x) <= tolerance and abs(first.y - second.y) <= tolerance


def _cubic_points(start, control_one, control_two, end, steps: int = 16) -> list[fitz.Point]:
    points = []
    for index in range(1, steps + 1):
        t = index / steps
        inverse = 1 - t
        points.append(
            fitz.Point(
                inverse ** 3 * start.x
                + 3 * inverse * inverse * t * control_one.x
                + 3 * inverse * t * t * control_two.x
                + t ** 3 * end.x,
                inverse ** 3 * start.y
                + 3 * inverse * inverse * t * control_one.y
                + 3 * inverse * t * t * control_two.y
                + t ** 3 * end.y,
            )
        )
    return points


def _drawing_contours(items: list[tuple]) -> list[list[fitz.Point]] | None:
    contours: list[list[fitz.Point]] = []

    def add_segment(start: fitz.Point, following: list[fitz.Point]) -> None:
        if not contours or not _points_match(contours[-1][-1], start):
            contours.append([fitz.Point(start)])
        contours[-1].extend(fitz.Point(point) for point in following)

    for item in items:
        kind = item[0]
        if kind == "l":
            add_segment(item[1], [item[2]])
        elif kind == "c":
            add_segment(item[1], _cubic_points(item[1], item[2], item[3], item[4]))
        elif kind == "re":
            rect = item[1]
            contours.append(
                [
                    fitz.Point(rect.x0, rect.y0),
                    fitz.Point(rect.x1, rect.y0),
                    fitz.Point(rect.x1, rect.y1),
                    fitz.Point(rect.x0, rect.y1),
                    fitz.Point(rect.x0, rect.y0),
                ]
            )
        elif kind == "qu":
            quad = item[1]
            contours.append([quad.ul, quad.ur, quad.lr, quad.ll, quad.ul])
        else:
            return None
    return [contour for contour in contours if len(contour) > 1]


def _add_drawing_fallback(slide, page, drawing: dict, to_emu_x, to_emu_y) -> None:
    rect = fitz.Rect(drawing.get("rect") or page.rect) & page.rect
    if rect.is_empty or rect.width <= 0 or rect.height <= 0:
        return
    pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), clip=rect, alpha=True)
    left, top = to_emu_x(rect.x0), to_emu_y(rect.y0)
    width, height = to_emu_x(rect.x1) - left, to_emu_y(rect.y1) - top
    if width > 0 and height > 0:
        slide.shapes.add_picture(BytesIO(pixmap.tobytes("png")), left, top, width=width, height=height)


def _add_powerpoint_drawing(slide, page, drawing: dict, scale: float, offset_x: float, offset_y: float, to_emu_x, to_emu_y) -> None:
    items = drawing.get("items") or []
    if len(items) == 1 and items[0][0] == "re":
        rect = items[0][1]
        left, top = to_emu_x(rect.x0), to_emu_y(rect.y0)
        width, height = to_emu_x(rect.x1) - left, to_emu_y(rect.y1) - top
        if width > 0 and height > 0:
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
            _style_powerpoint_shape(shape, drawing, scale)
        return

    if drawing.get("fill") is None and items and all(item[0] == "l" for item in items):
        for _, start, end in items:
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                to_emu_x(start.x),
                to_emu_y(start.y),
                to_emu_x(end.x),
                to_emu_y(end.y),
            )
            _style_powerpoint_shape(connector, drawing, scale, include_fill=False)
        return

    contours = _drawing_contours(items)
    if not contours:
        _add_drawing_fallback(slide, page, drawing, to_emu_x, to_emu_y)
        return

    coordinate_scale = EMU_PER_POINT * scale
    builder = slide.shapes.build_freeform(contours[0][0].x, contours[0][0].y, scale=coordinate_scale)
    for index, contour in enumerate(contours):
        if index:
            builder.move_to(contour[0].x, contour[0].y)
        closed = _points_match(contour[0], contour[-1]) or bool(drawing.get("closePath")) or drawing.get("fill") is not None
        vertices = contour[1:-1] if closed and _points_match(contour[0], contour[-1]) else contour[1:]
        builder.add_line_segments([(point.x, point.y) for point in vertices], close=closed)
    shape = builder.convert_to_shape(origin_x=Emu(int(offset_x)), origin_y=Emu(int(offset_y)))
    _style_powerpoint_shape(shape, drawing, scale)


def _add_powerpoint_image(slide, image: bytes, left: int, top: int, width: int, height: int) -> None:
    try:
        slide.shapes.add_picture(BytesIO(image), left, top, width=width, height=height)
        return
    except Exception:
        pass

    from PIL import Image

    picture = Image.open(BytesIO(image))
    if picture.mode not in ("RGB", "RGBA"):
        picture = picture.convert("RGBA" if "transparency" in picture.info else "RGB")
    normalized = BytesIO()
    picture.save(normalized, format="PNG")
    normalized.seek(0)
    slide.shapes.add_picture(normalized, left, top, width=width, height=height)


def _add_editable_slide(presentation, page, slide_w_emu: int, slide_h_emu: int) -> None:
    """Reconstruct a PDF page using editable text, images, and vector shapes.

    PDF drawing commands have no chart semantics, so charts become editable
    PowerPoint shapes rather than native chart objects. Unknown drawing commands
    are retained as clipped raster fallbacks instead of disappearing.
    """

    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)

    page_w = page.rect.width or 1
    page_h = page.rect.height or 1
    scale = min(slide_w_emu / (page_w * EMU_PER_POINT), slide_h_emu / (page_h * EMU_PER_POINT))
    offset_x = (slide_w_emu - page_w * EMU_PER_POINT * scale) / 2
    offset_y = (slide_h_emu - page_h * EMU_PER_POINT * scale) / 2

    def to_emu_x(value: float) -> int:
        return int(offset_x + value * EMU_PER_POINT * scale)

    def to_emu_y(value: float) -> int:
        return int(offset_y + value * EMU_PER_POINT * scale)

    data = page.get_text("dict")

    # Put raster content at the back, followed by vectors and editable text.
    for block in data["blocks"]:
        if block.get("type") != 1 or "image" not in block:
            continue
        x0, y0, x1, y1 = block["bbox"]
        width = to_emu_x(x1) - to_emu_x(x0)
        height = to_emu_y(y1) - to_emu_y(y0)
        if width <= 0 or height <= 0:
            continue
        _add_powerpoint_image(slide, block["image"], to_emu_x(x0), to_emu_y(y0), width, height)

    for drawing in page.get_drawings():
        _add_powerpoint_drawing(slide, page, drawing, scale, offset_x, offset_y, to_emu_x, to_emu_y)

    for block in data["blocks"]:
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            spans = [span for span in line.get("spans", []) if span.get("text", "").strip()]
            if not spans:
                continue
            x0, y0, x1, y1 = line["bbox"]
            left, top = to_emu_x(x0), to_emu_y(y0)
            width = max(to_emu_x(x1) - left, Emu(1))
            height = max(to_emu_y(y1) - top, Emu(1))
            textbox = slide.shapes.add_textbox(left, top, width, height)
            frame = textbox.text_frame
            frame.word_wrap = False
            frame.vertical_anchor = MSO_ANCHOR.TOP
            frame.margin_left = frame.margin_right = 0
            frame.margin_top = frame.margin_bottom = 0
            paragraph = frame.paragraphs[0]
            for span in spans:
                run = paragraph.add_run()
                run.text = span["text"]
                run.font.size = Pt(max(span.get("size", 12) * scale, 1))
                run.font.color.rgb = _pdf_int_color_to_rgb(int(span.get("color", 0)))
                flags = span.get("flags", 0)
                run.font.bold = bool(flags & 16)
                run.font.italic = bool(flags & 2)
                fontname = (span.get("font") or "").split("+")[-1]
                if fontname:
                    run.font.name = fontname


def pdf_to_powerpoint(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    editable = (options.get("pptx_mode") or "editable") != "image"
    document = fitz.open(source)
    try:
        presentation = Presentation()
        first = document[0].rect if len(document) else fitz.Rect(0, 0, 792, 612)
        presentation.slide_width = Emu(min(int(first.width * EMU_PER_POINT), Inches(56)))
        presentation.slide_height = Emu(min(int(first.height * EMU_PER_POINT), Inches(56)))
        if editable:
            for page in document:
                _add_editable_slide(presentation, page, presentation.slide_width, presentation.slide_height)
        else:
            dpi = int(options.get("dpi") or 160)
            blank_layout = presentation.slide_layouts[6]
            for page in document:
                slide = presentation.slides.add_slide(blank_layout)
                page_image = BytesIO(page.get_pixmap(dpi=dpi, alpha=False).tobytes("png"))
                slide.shapes.add_picture(page_image, 0, 0, width=presentation.slide_width, height=presentation.slide_height)
        presentation.save(output)
    finally:
        document.close()
    return output


WORD_EXTENSIONS = {".doc", ".docx", ".odt"}
EXCEL_EXTENSIONS = {".xls", ".xlsx", ".ods"}
POWERPOINT_EXTENSIONS = {".ppt", ".pptx", ".odp"}


def _office_to_pdf(files: list[Path], output: Path, allowed: set[str], label: str) -> Path:
    source = single_file(files)
    if source.suffix.lower() not in allowed:
        raise OperationError(f"{source.name} is not a {label} document")
    if os.name == "nt" and source.suffix.lower() in WORD_EXTENSIONS | EXCEL_EXTENSIONS | POWERPOINT_EXTENSIONS:
        try:
            return _windows_office_to_pdf(source, output)
        except (OperationError, subprocess.SubprocessError):
            pass
    soffice = resolve_binary("soffice", "libreoffice")
    subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output.parent),
            str(source),
        ],
        check=True,
    )
    generated = output.parent / f"{source.stem}.pdf"
    if not generated.exists():
        raise OperationError("LibreOffice did not produce a PDF")
    generated.replace(output)
    return output


def _windows_office_to_pdf(source: Path, output: Path) -> Path:
    """Use an installed Microsoft Office application as the Windows PDF backend."""
    powershell = shutil.which("powershell.exe") or str(
        Path(os.environ.get("SystemRoot", r"C:\Windows"))
        / "System32"
        / "WindowsPowerShell"
        / "v1.0"
        / "powershell.exe"
    )
    suffix = source.suffix.lower()
    if suffix in WORD_EXTENSIONS:
        body = (
            "$app=New-Object -ComObject Word.Application; $app.Visible=$false; "
            "try {$doc=$app.Documents.Open($src); $doc.ExportAsFixedFormat($dst,17)} "
            "finally {if($doc){$doc.Close($false)}; $app.Quit()}"
        )
    elif suffix in EXCEL_EXTENSIONS:
        body = (
            "$app=New-Object -ComObject Excel.Application; $app.Visible=$false; $app.DisplayAlerts=$false; "
            "try {$book=$app.Workbooks.Open($src); $book.ExportAsFixedFormat(0,$dst)} "
            "finally {if($book){$book.Close($false)}; $app.Quit()}"
        )
    elif suffix in POWERPOINT_EXTENSIONS:
        body = (
            "$app=New-Object -ComObject PowerPoint.Application; "
            "try {$deck=$app.Presentations.Open($src,$true,$false,$false); $deck.SaveAs($dst,32)} "
            "finally {if($deck){$deck.Close()}; $app.Quit()}"
        )
    else:
        raise OperationError(f"No Microsoft Office PDF backend for {source.suffix}")

    script = f"$src=$args[0]; $dst=$args[1]; $ErrorActionPreference='Stop'; {body}"
    subprocess.run(
        [powershell, "-NoProfile", "-NonInteractive", "-Command", script, str(source.resolve()), str(output.resolve())],
        check=True,
        timeout=120,
        capture_output=True,
        text=True,
    )
    if not output.exists():
        raise OperationError("Microsoft Office did not produce a PDF")
    return output


def office_to_pdf(files: list[Path], output: Path, options: dict) -> Path:
    return _office_to_pdf(files, output, OFFICE_EXTENSIONS, "supported Office")


def word_to_pdf(files: list[Path], output: Path, options: dict) -> Path:
    return _office_to_pdf(files, output, WORD_EXTENSIONS, "Word")


def excel_to_pdf(files: list[Path], output: Path, options: dict) -> Path:
    return _office_to_pdf(files, output, EXCEL_EXTENSIONS, "Excel")


def powerpoint_to_pdf(files: list[Path], output: Path, options: dict) -> Path:
    return _office_to_pdf(files, output, POWERPOINT_EXTENSIONS, "PowerPoint")


def _ocrmypdf_command() -> list[str]:
    """Prefer the standalone ocrmypdf executable, but fall back to running it via the
    current Python interpreter. The service process often lacks the venv's Scripts on
    PATH, so `-m ocrmypdf` guarantees the installed package is always reachable."""
    binary = shutil.which("ocrmypdf")
    if binary:
        return [binary]
    try:
        import ocrmypdf  # noqa: F401
    except ImportError as exc:
        raise OperationError(
            "OCR is not installed on the server. Install it with: pip install ocrmypdf"
        ) from exc
    return [sys.executable, "-m", "ocrmypdf"]


def _ocr_environment() -> dict:
    """ocrmypdf shells out to Tesseract and Ghostscript. Add their locations to PATH so
    OCR works even when the service process was not started from an activated shell."""
    env = dict(os.environ)
    extra_dirs = []
    for names in (("tesseract",), ("gswin64c", "gswin32c", "gs")):
        try:
            extra_dirs.append(str(Path(resolve_binary(*names)).parent))
        except OperationError:
            pass
    if extra_dirs:
        env["PATH"] = os.pathsep.join(extra_dirs + [env.get("PATH", "")])
    return env


def ocr_pdf(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    language = options.get("language") or "eng"
    command = _ocrmypdf_command()
    result = subprocess.run(
        [*command, "-l", language, "--skip-text", str(source), str(output)],
        env=_ocr_environment(),
        capture_output=True,
        text=True,
    )
    if result.returncode != 0:
        message = (result.stderr or result.stdout or "").strip()
        lowered = message.lower()
        if "tesseract" in lowered:
            raise OperationError(
                "OCR engine (Tesseract) is not installed or not reachable on the server. "
                "Install Tesseract-OCR and retry."
            )
        if "ghostscript" in lowered or "gswin" in lowered:
            raise OperationError(
                "Ghostscript is not installed or not reachable on the server. Install Ghostscript and retry."
            )
        last_line = message.splitlines()[-1] if message else "unknown error"
        raise OperationError(f"OCR failed: {last_line}")
    return output


def flatten_edits(files: list[Path], output: Path, options: dict) -> Path:
    source = single_pdf(files)
    annotations = options.get("annotations") or {}
    pages = annotations.get("pages") if isinstance(annotations, dict) else None
    if not isinstance(pages, list):
        raise OperationError("Invalid annotation document")

    signature_fields: list[tuple[int, dict]] = []
    document = fitz.open(source)
    try:
        for page_data in pages:
            page_index = int(page_data.get("page", 0))
            if page_index < 0 or page_index >= len(document):
                raise OperationError("Annotation references a missing page")
            page = document[page_index]
            objects = page_data.get("objects") or []

            for item in objects:
                if item.get("type") == "text" and item.get("erase", True):
                    page.add_redact_annot(normalized_erase_rect(page, item), fill=(1, 1, 1))
            if any(item.get("type") == "text" and item.get("erase", True) for item in objects):
                page.apply_redactions()

            for item in objects:
                if item.get("type") == "signature_field":
                    signature_fields.append((page_index, item))
                draw_annotation(page, item)
        document.save(output, garbage=4, deflate=True)
    finally:
        document.close()
    if signature_fields:
        add_signature_widgets(output, signature_fields)
    return output


def draw_annotation(page, item: dict) -> None:
    kind = item.get("type")
    color = html_color(item.get("color", "#111827"))
    opacity = min(max(float(item.get("opacity", 1)), 0), 1)
    if kind == "text":
        rect = normalized_rect(page, item)
        fontsize = max(float(item.get("font_size", 16)) * page.rect.width / float(item.get("viewport_width") or page.rect.width), 6)
        rect = fitz.Rect(rect.x0, rect.y0, rect.x1, min(max(rect.y1, rect.y0 + fontsize * 1.4), page.rect.y1))
        family = str(item.get("font_family") or "helvetica").lower()
        family_key = "times" if "times" in family else "courier" if "courier" in family else "helvetica"
        fontname = PAGE_NUMBER_FONTS.get((family_key, _truthy(item.get("bold")), _truthy(item.get("italic"))), "helv")
        align = {"center": fitz.TEXT_ALIGN_CENTER, "right": fitz.TEXT_ALIGN_RIGHT}.get(item.get("align"), fitz.TEXT_ALIGN_LEFT)
        text = str(item.get("text", ""))
        text_width = min(fitz.get_text_length(text, fontname=fontname, fontsize=fontsize), max(page.rect.x1 - rect.x0, 1))
        x0 = rect.x0 if align == fitz.TEXT_ALIGN_LEFT else rect.x1 - text_width if align == fitz.TEXT_ALIGN_RIGHT else rect.x0 + (rect.width - text_width) / 2
        if text and "\n" not in text:
            page.insert_text(fitz.Point(max(x0, 0), min(rect.y0 + fontsize, page.rect.y1)), text, fontsize=fontsize, fontname=fontname, color=color)
        elif text:
            page.insert_textbox(rect, text, fontsize=fontsize, fontname=fontname, color=color, align=align)
        if _truthy(item.get("underline")):
            y = min(rect.y0 + fontsize * 1.15, rect.y1)
            page.draw_line(fitz.Point(x0, y), fitz.Point(x0 + text_width, y), color=color, width=max(fontsize * 0.045, 0.6))
    elif kind in {"rectangle", "highlight"}:
        rect = normalized_rect(page, item)
        fill = html_color(item.get("fill", "#fde047")) if item.get("fill") else None
        page.draw_rect(
            rect,
            color=color if kind == "rectangle" else None,
            fill=fill,
            width=max(float(item.get("stroke_width", 2)), 0.5),
            stroke_opacity=opacity,
            fill_opacity=opacity,
        )
    elif kind == "circle":
        rect = normalized_rect(page, item)
        fill = html_color(item.get("fill")) if item.get("fill") else None
        page.draw_oval(
            rect,
            color=color,
            fill=fill,
            width=max(float(item.get("stroke_width", 2)), 0.5),
            stroke_opacity=opacity,
            fill_opacity=opacity,
        )
    elif kind == "path":
        points = item.get("points") or []
        if len(points) > 1:
            shape = page.new_shape()
            for first, second in zip(points, points[1:]):
                shape.draw_line(normalized_point(page, first), normalized_point(page, second))
            shape.finish(color=color, width=max(float(item.get("stroke_width", 2)), 0.5), stroke_opacity=opacity)
            shape.commit()
    elif kind in {"line", "arrow"}:
        rect = normalized_rect(page, item)
        start = fitz.Point(rect.x0, rect.y0)
        end = fitz.Point(rect.x1, rect.y1)
        width = max(float(item.get("stroke_width", 2)), 0.5)
        page.draw_line(start, end, color=color, width=width, stroke_opacity=opacity)
        if kind == "arrow":
            angle = math.atan2(end.y - start.y, end.x - start.x)
            head = max(min(rect.width, rect.height, 22), 9)
            for offset in (-0.55, 0.55):
                point = fitz.Point(end.x - head * math.cos(angle + offset), end.y - head * math.sin(angle + offset))
                page.draw_line(end, point, color=color, width=width, stroke_opacity=opacity)
    elif kind == "image":
        image_bytes = decode_data_url(item.get("image", ""))
        if image_bytes:
            page.insert_image(normalized_rect(page, item), stream=image_bytes, keep_proportion=True, overlay=True)
    elif kind == "stamp":
        rect = normalized_rect(page, item)
        text = str(item.get("text") or "APPROVED")
        page.draw_rect(rect, color=color, width=max(float(item.get("stroke_width", 2)), 1), stroke_opacity=opacity)
        fontsize = max(min(rect.height * 0.45, rect.width / max(len(text) * 0.62, 1), 28), 6)
        text_width = fitz.get_text_length(text, fontname="hebo", fontsize=fontsize)
        origin = fitz.Point(rect.x0 + max((rect.width - text_width) / 2, 1), rect.y0 + (rect.height + fontsize * 0.72) / 2)
        page.insert_text(origin, text, fontsize=fontsize, fontname="hebo", color=color)
    elif kind == "signature_field":
        rect = normalized_rect(page, item)
        page.draw_rect(rect, color=(0.12, 0.37, 0.60), fill=(0.84, 0.91, 0.97), width=1, fill_opacity=0.55)
        page.insert_textbox(rect, "Sign here", fontsize=max(min(rect.height * 0.34, 14), 7), fontname="hebo", color=(0.09, 0.23, 0.45), align=fitz.TEXT_ALIGN_CENTER)
    elif kind == "signature":
        rect = normalized_rect(page, item)
        image_bytes = decode_data_url(item.get("image", ""))
        if image_bytes:
            page.insert_image(rect, stream=image_bytes, keep_proportion=True, overlay=True)
        signer = str(item.get("signer_name", "")).strip()
        signed_at = str(item.get("signed_at", "")).strip()
        if signer:
            stamp = fitz.Rect(rect.x0, rect.y1 + 2, rect.x1, min(rect.y1 + 24, page.rect.height))
            page.insert_textbox(stamp, f"Signed by {signer} on {signed_at}", fontsize=7, fontname="helv", color=(0.18, 0.22, 0.28))


def normalized_rect(page, item: dict):
    x = float(item.get("x", 0)) * page.rect.width
    y = float(item.get("y", 0)) * page.rect.height
    width = max(float(item.get("width", 0)) * page.rect.width, 1)
    height = max(float(item.get("height", 0)) * page.rect.height, 1)
    return fitz.Rect(x, y, min(x + width, page.rect.width), min(y + height, page.rect.height))


def normalized_erase_rect(page, item: dict):
    if not all(key in item for key in ("erase_x", "erase_y", "erase_width", "erase_height")):
        return normalized_rect(page, item)
    erase_item = {
        "x": item["erase_x"],
        "y": item["erase_y"],
        "width": item["erase_width"],
        "height": item["erase_height"],
    }
    return normalized_rect(page, erase_item)


def normalized_point(page, point):
    return fitz.Point(float(point[0]) * page.rect.width, float(point[1]) * page.rect.height)


def add_signature_widgets(path: Path, fields: list[tuple[int, dict]]) -> None:
    reader = PdfReader(str(path))
    writer = PdfWriter()
    writer.clone_document_from_reader(reader)
    root = writer._root_object
    acroform = root.get("/AcroForm")
    if acroform is None:
        acroform = DictionaryObject({NameObject("/Fields"): ArrayObject(), NameObject("/SigFlags"): NumberObject(3)})
        root[NameObject("/AcroForm")] = writer._add_object(acroform)
    else:
        acroform = acroform.get_object()
        if "/Fields" not in acroform:
            acroform[NameObject("/Fields")] = ArrayObject()
        acroform[NameObject("/SigFlags")] = NumberObject(3)
    form_fields = acroform["/Fields"]

    for page_index, item in fields:
        page = writer.pages[page_index]
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        x0 = float(item.get("x", 0)) * width
        y_top = float(item.get("y", 0)) * height
        x1 = min(x0 + max(float(item.get("width", 0)) * width, 1), width)
        y_bottom = min(y_top + max(float(item.get("height", 0)) * height, 1), height)
        field = DictionaryObject({
            NameObject("/Type"): NameObject("/Annot"),
            NameObject("/Subtype"): NameObject("/Widget"),
            NameObject("/FT"): NameObject("/Sig"),
            NameObject("/T"): TextStringObject(str(item.get("field_name") or "SignatureFormField")[:200]),
            NameObject("/Rect"): ArrayObject([FloatObject(x0), FloatObject(height - y_bottom), FloatObject(x1), FloatObject(height - y_top)]),
            NameObject("/F"): NumberObject(4),
            NameObject("/P"): page.indirect_reference,
        })
        field_ref = writer._add_object(field)
        annotations = page.get("/Annots")
        if annotations is None:
            annotations = ArrayObject()
            page[NameObject("/Annots")] = annotations
        else:
            annotations = annotations.get_object()
        annotations.append(field_ref)
        form_fields.append(field_ref)

    temporary = path.with_name(f"{path.stem}.widgets{path.suffix}")
    with temporary.open("wb") as handle:
        writer.write(handle)
    temporary.replace(path)


def html_color(value: str):
    value = value.lstrip("#")
    if len(value) != 6:
        return (0.07, 0.09, 0.13)
    try:
        return tuple(int(value[index:index + 2], 16) / 255 for index in (0, 2, 4))
    except ValueError:
        return (0.07, 0.09, 0.13)


def decode_data_url(value: str) -> bytes | None:
    if not value.startswith("data:image/") or "," not in value:
        return None
    try:
        return base64.b64decode(value.split(",", 1)[1], validate=True)
    except (ValueError, TypeError):
        return None


def resolve_binary(*names: str) -> str:
    for name in names:
        path = shutil.which(name)
        if path:
            return path
        if os.name == "nt":
            for extension in (".com", ".exe", ".bat", ".cmd"):
                path = shutil.which(f"{name}{extension}")
                if path:
                    return path

    if os.name == "nt":
        local_tools = Path(__file__).resolve().parent.parent / ".tools"
        roots = [
            Path(os.environ.get("ProgramFiles", r"C:\Program Files")),
            Path(os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)")),
        ]
        candidates = [
            local_tools / "libreoffice" / "program" / "soffice.com",
            local_tools / "libreoffice" / "program" / "soffice.exe",
            local_tools / "ghostscript" / "bin" / "gswin64c.exe",
            local_tools / "tesseract" / "tesseract.exe",
        ]
        for root in roots:
            candidates.extend(
                [
                    root / "LibreOffice" / "program" / "soffice.com",
                    root / "LibreOffice" / "program" / "soffice.exe",
                    root / "Tesseract-OCR" / "tesseract.exe",
                ]
            )
            candidates.extend(root.glob("gs/gs*/bin/gswin64c.exe"))
            candidates.extend(root.glob("gs/gs*/bin/gswin32c.exe"))
        requested = {name.lower() for name in names}
        for candidate in candidates:
            if candidate.exists() and candidate.stem.lower() in requested:
                return str(candidate)
    raise OperationError(f"Missing system dependency: one of {', '.join(names)}")


def ensure_pdf(path: Path) -> None:
    if path.suffix.lower() not in PDF_EXTENSIONS:
        raise OperationError(f"{path.name} is not a PDF")


def single_file(files: list[Path]) -> Path:
    if len(files) != 1:
        raise OperationError("This operation requires exactly one file")
    return files[0]


def single_pdf(files: list[Path]) -> Path:
    source = single_file(files)
    ensure_pdf(source)
    return source


def parse_page_spec(spec: str, page_count: int) -> list[int]:
    pages: list[int] = []
    for raw_part in spec.replace(" ", "").split(","):
        if not raw_part:
            continue
        if "-" in raw_part:
            start_text, end_text = raw_part.split("-", 1)
            start = int(start_text)
            end = int(end_text)
            step = 1 if end >= start else -1
            pages.extend(range(start, end + step, step))
        else:
            pages.append(int(raw_part))
    for page_number in pages:
        if page_number < 1 or page_number > page_count:
            raise OperationError(f"Page {page_number} is outside the document")
    return pages


OPERATIONS = {
    "edit": ("edited.pdf", flatten_edits),
    "merge": ("merged.pdf", merge_pdfs),
    "split": ("split_pages.zip", split_pdf),
    "delete_pages": ("pages_deleted.pdf", delete_pages),
    "reorder_pages": ("reordered.pdf", reorder_pages),
    "rotate": ("rotated.pdf", rotate_pdf),
    "crop": ("cropped.pdf", crop_pdf),
    "compress": ("compressed.pdf", compress_pdf),
    "protect": ("protected.pdf", protect_pdf),
    "unlock": ("unlocked.pdf", unlock_pdf),
    "watermark": ("watermarked.pdf", watermark_pdf),
    "remove_watermark": ("watermark_removed.pdf", remove_watermark_pdf),
    "page_numbers": ("page_numbers.pdf", page_numbers_pdf),
    "extract_images": ("extracted_images.zip", extract_images),
    "images_to_pdf": ("images.pdf", images_to_pdf),
    "pdf_to_images": ("pdf_images.zip", pdf_to_images),
    "pdf_to_word": ("converted.docx", pdf_to_word),
    "pdf_to_excel": ("converted.xlsx", pdf_to_excel),
    "word_to_excel": ("converted.xlsx", word_to_excel),
    "pdf_to_powerpoint": ("converted.pptx", pdf_to_powerpoint),
    "office_to_pdf": ("converted.pdf", office_to_pdf),
    "word_to_pdf": ("converted.pdf", word_to_pdf),
    "excel_to_pdf": ("converted.pdf", excel_to_pdf),
    "powerpoint_to_pdf": ("converted.pdf", powerpoint_to_pdf),
    "html_to_pdf": ("webpage.pdf", html_to_pdf),
    "ocr": ("searchable.pdf", ocr_pdf),
}
